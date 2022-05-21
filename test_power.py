from pptx import Presentation

prs = Presentation()
#title_slide_layout = prs.slide_layouts[0]
title_slide_layout = prs.slide_layouts[1] 


slide = prs.slides.add_slide(title_slide_layout)

#追加：シートの追加
slide2 = prs.slides.add_slide(title_slide_layout)

title = slide.shapes.title

#追加:
title2 = slide2.shapes.title
subtitle = slide2.placeholders[1] 

subtitle = slide.placeholders[1]

title.text = "Hello, World!"
title2.text = "Welcome"

subtitle.text = "python-pptx was here!"
subtitle.text = "!World!"
#prs.save('test.pptx')
prs.save('test2.pptx')    
